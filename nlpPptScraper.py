from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import os
import logging
import re
import spacy 
import json
from spacy_lookup import Entity
from spacy.pipeline import EntityRuler

nlp = spacy.blank('en')

def main():  
    # Root folder for audit
    root = ""
    # Folder where ppts are located
    pptFolder = ""
    # Absolute Path where the files are located
    targetDir = f"{root}/{pptFolder}"
    # Import List of precedents in Excel. 
    # Can be substituted for a simple list of phrases
    listofValues = getExcelColumn(f"{root}/ref_data.xlsx",'Actor')
    # TODO: Make this a separate function
    # Initialize NLP pipeline
    patterns = []
    for value in listofValues:
        patterns.append({"label": "ORG", "pattern": value})
    ruler = EntityRuler(nlp, overwrite_ents=True)
    ruler.add_patterns(patterns)
    nlp.add_pipe(ruler)
    phrase_matcher = ""
    # Find all powerpoints in target directory
    files = [x for x in os.listdir(targetDir) if x.endswith(".pptx")] 
    # Make a new folder for parsed output
    newDirName = "parsed"
    newDir = os.path.join(targetDir, newDirName) 
    os.mkdir(newDir) 
    # Run parser
    pptScraper(files, phrase_matcher, targetDir, newDir)

def pptScraper(files, phrase_matcher, oldDir, newDir, *args):
    # TODO: Add a default flag and initialize in main
    txtFlag = "*ALREADY IN DATABASE*"
    for file in files:
        f = open(f"{oldDir}/{file}", "rb")
        prs = Presentation(f) 
        currentPres = str(prs.core_properties.title)
        print("Running on " + file)
        # Parse shapes on each slide
        for slide in prs.slides:
            # Create an aggregate string for all the words in the slide
            slideText = ""
            # Get text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slideText += shape.text
            doc = nlp(slideText)
            # If the key phrase exists in the slide, add a text or image flag
            if len([(ent.text, ent.label_) for ent in doc.ents]) > 0:
                addText(txtFlag, slide)
        # Saved the flagged file to the new directory
        newFile = f"{newDir}/parsed_{file}.pptx"
        saveFile(prs, newFile)
    return

def addImage(img, slide):
    left = Inches(5)
    top = Inches(0)
    slide.shapes.add_picture(img, left, top)
    return

def addText(textFlag, slide):
    left = Inches(5)
    top = Inches(0)
    width = Inches(5)
    height = Inches(0.5)
    color = RGBColor(255, 255, 0)
    newShape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    # Add Text
    run = newShape.text_frame.paragraphs[0].add_run()
    run.font.color.rgb = RGBColor(0, 0, 0)
    run.text = textFlag
    # Add outline
    shapeLine = newShape.line
    shapeLine.color.rgb = color
    # Add a fill
    shapeFill = newShape.fill
    shapeFill.solid()
    shapeFill.fore_color.rgb = color
    return

def getExcelColumn(filepath, col_name):
    df = pd.read_excel(filepath)
    colToList = df[col_name].dropna().tolist()
    return colToList

def saveFile(prs, path):
    prs.save(path)
    print(f"Saved to {path}")
    return

main()
